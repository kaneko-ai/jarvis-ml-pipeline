"""PPTX export helpers for writing drafts."""

from __future__ import annotations

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:  # pragma: no cover
    PPTX_AVAILABLE = False
    Presentation = None
    Inches = None
    Pt = None

from ..writing.outline_builder import ClaimDatum
from ..writing.utils import load_overview

SLIDE_WIDTH = Inches(10) if Inches else None
SLIDE_HEIGHT = Inches(7.5) if Inches else None


def _slide_footer(slide, text: str) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.4))
    text_frame = textbox.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)


def _add_bullets(presentation: Presentation, title: str, bullets: list[str]):
    layout = presentation.slide_layouts[1]
    content_slide = presentation.slides.add_slide(layout)
    content_slide.shapes.title.text = title
    body = content_slide.shapes.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
    return content_slide


def _evidence_footer_text(claim: ClaimDatum | None) -> str:
    if not claim or not claim.evidence:
        return "Evidence: unknown"
    ev = claim.evidence[0]
    return f"Evidence: {ev.paper_id}/{ev.chunk_id}"


def build_pptx_from_slides(run_dir: Path, claims: list[ClaimDatum], output_path: Path) -> Path:
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed. Please add python-pptx to requirements.")

    presentation = Presentation()
    overview = load_overview(run_dir)
    ordered = claims
    top_claim = ordered[0] if ordered else None
    second_claim = ordered[1] if len(ordered) > 1 else None

    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "研究テーマ"
    slide.placeholders[1].text = "所属 / 氏名"
    _slide_footer(slide, _evidence_footer_text(top_claim))

    background_bullets = [
        (overview.split("\n")[0].strip() if overview else "Background summary placeholder."),
        "図プレースホルダ: 概念図を挿入してください。",
    ]
    background_slide = _add_bullets(presentation, "Background", background_bullets)
    _slide_footer(background_slide, _evidence_footer_text(top_claim))

    gap_bullets = [
        f"Fact: {second_claim.text if second_claim else 'Gap placeholder'}",
        "Interpretation: 既存研究の不足点を整理する。",
    ]
    gap_slide = _add_bullets(presentation, "Gap & Hypothesis", gap_bullets)
    _slide_footer(gap_slide, _evidence_footer_text(second_claim or top_claim))

    aims_bullets = [
        f"Aim 1: {top_claim.text if top_claim else 'Aim placeholder'}",
        "Aim 2: [[YOUR_DATA_HERE]]",
        "Aim 3: [[YOUR_DATA_HERE]]",
    ]
    aims_slide = _add_bullets(presentation, "Aim / Objective", aims_bullets)
    _slide_footer(aims_slide, _evidence_footer_text(top_claim))

    methods_bullets = [
        "Fact: 既存の方法概要を記述。",
        "Interpretation: 方法の意図と利点を説明。",
        "模式図プレースホルダ: 実験フローを挿入。",
    ]
    methods_slide = _add_bullets(presentation, "Methods overview", methods_bullets)
    _slide_footer(methods_slide, _evidence_footer_text(top_claim))

    result1_claim = ordered[0] if ordered else None
    result2_claim = ordered[1] if len(ordered) > 1 else None

    result1_bullets = [
        f"Fact: {result1_claim.text if result1_claim else 'Result placeholder'}",
        "Interpretation: 結果の示唆を記述。",
        "図プレースホルダ: 結果図を挿入。",
    ]
    result1_slide = _add_bullets(presentation, "Key Result 1", result1_bullets)
    _slide_footer(result1_slide, _evidence_footer_text(result1_claim))

    result2_bullets = [
        f"Fact: {result2_claim.text if result2_claim else 'Result placeholder'}",
        "Interpretation: 結果の示唆を記述。",
        "図プレースホルダ: 結果図を挿入。",
    ]
    result2_slide = _add_bullets(presentation, "Key Result 2", result2_bullets)
    _slide_footer(result2_slide, _evidence_footer_text(result2_claim))

    discussion_bullets = [
        "Fact: 結果の再掲。",
        "Interpretation: 限界と今後の課題を整理。",
        "Limitations: [[LAB_TERMS_CHECK]]",
    ]
    discussion_slide = _add_bullets(presentation, "Discussion", discussion_bullets)
    _slide_footer(discussion_slide, _evidence_footer_text(top_claim))

    future_bullets = [
        "Next experiment: [[YOUR_DATA_HERE]]",
        "New hypothesis: [[YOUR_DATA_HERE]]",
    ]
    future_slide = _add_bullets(presentation, "Future work", future_bullets)
    _slide_footer(future_slide, _evidence_footer_text(top_claim))

    takehome_bullets = [
        f"Fact: {top_claim.text if top_claim else 'Take-home placeholder'}",
        "Interpretation: 研究の主要な学術的意義。",
    ]
    takehome_slide = _add_bullets(presentation, "Take-home message", takehome_bullets)
    _slide_footer(takehome_slide, _evidence_footer_text(top_claim))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return output_path