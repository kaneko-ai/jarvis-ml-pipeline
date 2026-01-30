"""Build QA reports for lab submission readiness."""

from __future__ import annotations

import importlib.util
import json
from pathlib import Path

from .checklists import run_checklists
from .figure_table_registry import scan_text
from .scientific_linter import ScientificLinter
from .term_normalizer import (
    load_style_guide as _load_style_guide,
)
from .term_normalizer import (
    normalize_docx_paragraphs,
    normalize_markdown,
    normalize_pptx_slides,
)

STYLE_GUIDE_PATH = Path(__file__).with_name("lab_style_guide.yaml")

DOCX_AVAILABLE = importlib.util.find_spec("docx") is not None
PPTX_AVAILABLE = importlib.util.find_spec("pptx") is not None

if DOCX_AVAILABLE:
    import docx

if PPTX_AVAILABLE:
    import pptx


def load_style_guide(path: Path = STYLE_GUIDE_PATH) -> dict[str, object]:
    return _load_style_guide(path)


def _load_docx_text(docx_path: Path) -> list[str]:
    if not DOCX_AVAILABLE:
        return []
    document = docx.Document(str(docx_path))
    return [p.text for p in document.paragraphs]


def _load_pptx_text(pptx_path: Path) -> list[str]:
    if not PPTX_AVAILABLE:
        return []
    presentation = pptx.Presentation(str(pptx_path))
    slides_text: list[str] = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_text.append("\n".join([t for t in slide_text if t]))
    return slides_text


def _count_by_issue_type(issues: list[dict[str, object]]) -> dict[str, int]:
    counts: dict[str, int] = {}
    for issue in issues:
        issue_type = issue.get("issue_type") or issue.get("type")
        if not issue_type:
            continue
        counts[issue_type] = counts.get(issue_type, 0) + 1
    return counts


def _summarize_conclusion(text: str) -> int:
    for section in ["Conclusion", "結論", "まとめ"]:
        if section in text:
            segment = text.split(section, 1)[-1]
            sentences = [s for s in segment.split("。") if s.strip()]
            return len(sentences)
    return 0


def run_qa_gate(
    run_id: str,
    run_dir: Path,
    targets: list[str] | None = None,
    output_base: Path = Path("data/runs"),
) -> dict[str, object]:
    """Run QA gate and generate reports."""
    style_guide = load_style_guide()
    targets = targets or ["md", "docx", "pptx"]

    issues: list[dict[str, object]] = []
    normalized: dict[str, object] = {}
    replacements: list[dict[str, str]] = []
    lint_issues: list[dict[str, object]] = []

    md_text = ""
    if "md" in targets:
        report_path = run_dir / "report.md"
        if report_path.exists():
            md_text = report_path.read_text(encoding="utf-8")
            normalized_md, md_issues, md_replacements = normalize_markdown(md_text, style_guide)
            normalized["md"] = normalized_md
            issues.extend([issue.__dict__ for issue in md_issues])
            replacements.extend(md_replacements)
            lint_issues.extend(
                [issue.__dict__ for issue in ScientificLinter().lint_text(normalized_md)]
            )

    if "docx" in targets:
        for docx_path in run_dir.glob("*.docx"):
            paragraphs = _load_docx_text(docx_path)
            if not paragraphs:
                continue
            result = normalize_docx_paragraphs(paragraphs, style_guide)
            normalized[docx_path.name] = "\n".join(result.normalized_lines)
            issues.extend([issue.__dict__ for issue in result.issues])
            replacements.extend(result.replacements)
            lint_issues.extend(
                [
                    issue.__dict__
                    for issue in ScientificLinter().lint_text("\n".join(result.normalized_lines))
                ]
            )

    if "pptx" in targets:
        for pptx_path in run_dir.glob("*.pptx"):
            slides = _load_pptx_text(pptx_path)
            if not slides:
                continue
            result = normalize_pptx_slides(slides, style_guide)
            normalized[pptx_path.name] = "\n".join(result.normalized_lines)
            issues.extend([issue.__dict__ for issue in result.issues])
            replacements.extend(result.replacements)
            lint_issues.extend(
                [
                    issue.__dict__
                    for issue in ScientificLinter().lint_text("\n".join(result.normalized_lines))
                ]
            )

    issues.extend(lint_issues)

    registry = scan_text(md_text) if md_text else {"figures": [], "tables": [], "issues": []}
    issues.extend(
        [
            {
                "issue_type": issue["issue_type"],
                "severity": "WARN" if issue["issue_type"] == "missing_reference" else "ERROR",
                "message": issue["message"],
                "location": "figure_registry",
                "excerpt": issue["fig_id"],
            }
            for issue in registry.get("issues", [])
        ]
    )

    error_count = sum(1 for issue in issues if issue.get("severity") == "ERROR")
    warn_count = sum(1 for issue in issues if issue.get("severity") == "WARN")
    info_count = sum(1 for issue in issues if issue.get("severity") == "INFO")
    ready_to_submit = error_count == 0

    counts = _count_by_issue_type(issues)
    conclusion_sentences = _summarize_conclusion(md_text) if md_text else 0

    qa_result: dict[str, object] = {
        "run_id": run_id,
        "ready_to_submit": ready_to_submit,
        "error_count": error_count,
        "warn_count": warn_count,
        "info_count": info_count,
        "issues": issues,
        "registry": registry,
        "normalized_text": normalized,
        "replacements": replacements,
        "counts": counts,
        "conclusion_sentences": conclusion_sentences,
    }

    qa_result["checklists"] = run_checklists(qa_result)
    qa_result["top_errors"] = [
        {
            "issue_type": issue.get("issue_type"),
            "message": issue.get("message"),
            "location": issue.get("location"),
        }
        for issue in issues
        if issue.get("severity") == "ERROR"
    ][:5]

    output_dir = output_base / run_id / "qa"
    output_dir.mkdir(parents=True, exist_ok=True)

    report_md = _build_markdown_report(qa_result)
    report_json_path = output_dir / "qa_report.json"
    report_md_path = output_dir / "qa_report.md"

    with open(report_json_path, "w", encoding="utf-8") as f:
        json.dump(qa_result, f, indent=2, ensure_ascii=False)

    report_md_path.write_text(report_md, encoding="utf-8")

    report_html_path = output_dir / "qa_report.html"
    report_html_path.write_text(_build_html_report(qa_result), encoding="utf-8")

    _sync_to_run_dir(run_dir, output_dir)

    return qa_result


def _build_markdown_report(qa_result: dict[str, object]) -> str:
    lines = [
        f"# QA Report: {qa_result.get('run_id')}",
        "",
        f"READY_TO_SUBMIT: {'true' if qa_result.get('ready_to_submit') else 'false'}",
        "",
        f"ERROR: {qa_result.get('error_count', 0)} / WARN: {qa_result.get('warn_count', 0)}",
        "",
        "## Top Errors",
    ]
    top_errors = qa_result.get("top_errors", [])
    if not top_errors:
        lines.append("- None")
    else:
        for err in top_errors:
            lines.append(f"- {err.get('issue_type')}: {err.get('message')} ({err.get('location')})")

    lines.append("")
    lines.append("## Issues")
    for issue in qa_result.get("issues", []):
        severity = issue.get("severity")
        issue_type = issue.get("issue_type")
        message = issue.get("message") or issue.get("original")
        location = issue.get("location")
        lines.append(f"- [{severity}] {issue_type}: {message} ({location})")

    lines.append("")
    lines.append("## Checklists")
    for item in qa_result.get("checklists", []):
        status = "✅" if item.get("passed") else "❌"
        lines.append(f"- {status} {item.get('description')}")

    return "\n".join(lines)


def _build_html_report(qa_result: dict[str, object]) -> str:
    title = f"QA Report: {qa_result.get('run_id')}"
    rows = "".join(
        f"<li>[{issue.get('severity')}] {issue.get('issue_type')}: {issue.get('message') or issue.get('original')}"
        f" <em>{issue.get('location')}</em></li>"
        for issue in qa_result.get("issues", [])
    )
    checklist_rows = "".join(
        f"<li>{'✅' if item.get('passed') else '❌'} {item.get('description')}</li>"
        for item in qa_result.get("checklists", [])
    )
    return (
        "<html><head><meta charset='utf-8'><title>" + title + "</title></head><body>"
        f"<h1>{title}</h1>"
        f"<p>READY_TO_SUBMIT: {'true' if qa_result.get('ready_to_submit') else 'false'}</p>"
        f"<p>ERROR: {qa_result.get('error_count', 0)} / WARN: {qa_result.get('warn_count', 0)}</p>"
        "<h2>Issues</h2><ul>" + rows + "</ul>"
        "<h2>Checklists</h2><ul>" + checklist_rows + "</ul>"
        "</body></html>"
    )


def _sync_to_run_dir(run_dir: Path, output_dir: Path) -> None:
    qa_dir = run_dir / "qa"
    qa_dir.mkdir(parents=True, exist_ok=True)
    for filename in ["qa_report.json", "qa_report.md", "qa_report.html"]:
        src = output_dir / filename
        if src.exists():
            dest = qa_dir / filename
            dest.write_bytes(src.read_bytes())
